# generate_presentation.py
"""
Generate a polished PowerPoint deck (NextBestAction.pptx) that summarises the
Standard Bank Next Best Action project.

The deck includes:
1. Title slide with project name and logo placeholder.
2. Problem statement (pulled from docs/problem_statement.md if present).
3. Classification Engine – data description, model, evaluation metrics.
4. Recommendation Engine – dummy recommendation table snapshot.
5. Next steps & conclusions.

The script uses `python-pptx` and can be run after the classification script
has produced `model/metrics.json` and after the dummy recommendation CSV is
available.
"""

import pathlib
import json
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# ---------------------------------------------------------------------------
# Configuration – adjust paths if you change the folder layout
# ---------------------------------------------------------------------------
BASE_DIR = pathlib.Path("..")  # Repository root (script lives in scripts/)
METRICS_PATH = BASE_DIR / "model" / "metrics.json"
RECOMMENDATIONS_PATH = BASE_DIR / "data" / "recommendations" / "dummy_recommendations.csv"
OUTPUT_PPTX = BASE_DIR / "NextBestAction.pptx"

# ---------------------------------------------------------------------------
# Helper to add a title slide
# ---------------------------------------------------------------------------
def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide_layout = prs.slide_layouts[0]  # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle

# ---------------------------------------------------------------------------
# Helper to add a bullet‑point slide
# ---------------------------------------------------------------------------
def add_bullet_slide(prs: Presentation, heading: str, bullets: list) -> None:
    slide_layout = prs.slide_layouts[1]  # Title & Content
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = heading
    tf = slide.shapes.placeholders[1].text_frame
    for i, bullet in enumerate(bullets):
        p = tf.add_paragraph() if i else tf.paragraphs[0]
        p.text = bullet
        p.level = 0
        p.font.size = Pt(18)

# ---------------------------------------------------------------------------
# Helper to add a table slide (used for recommendation snapshot)
# ---------------------------------------------------------------------------
def add_table_slide(prs: Presentation, heading: str, df: pd.DataFrame, max_rows: int = 8) -> None:
    slide_layout = prs.slide_layouts[5]  # Title Only
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = heading
    # Trim rows for readability
    display_df = df.head(max_rows)
    rows, cols = display_df.shape
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(0.5 + rows * 0.3)
    table = slide.shapes.add_table(rows + 1, cols, left, top, width, height).table
    # Header
    for col_idx, col_name in enumerate(display_df.columns):
        cell = table.cell(0, col_idx)
        cell.text = str(col_name)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(14)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0xD9, 0xD9, 0xD9)
    # Body rows
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r + 1, c)
            cell.text = str(display_df.iat[r, c])
            cell.text_frame.paragraphs[0].font.size = Pt(12)

# ---------------------------------------------------------------------------
# Main generation routine
# ---------------------------------------------------------------------------
def main() -> None:
    prs = Presentation()

    # 1️⃣ Title slide
    add_title_slide(prs, "Standard Bank – Next Best Action", "Classification + Recommendation Engine")

    # 2️⃣ Problem statement (static text – can be replaced with a markdown import)
    problem_bullets = [
        "Predict churn vs. non‑churn customers using a high‑performing classifier.",
        "Provide product recommendations to convert churned customers into loyal ones.",
        "Goal: improve conversion rate and long‑term relationship value.",
    ]
    add_bullet_slide(prs, "Project Overview", problem_bullets)

    # 3️⃣ Classification Engine – Metrics
    if METRICS_PATH.exists():
        with METRICS_PATH.open() as f:
            metrics = json.load(f)
        metric_bullets = [
            f"Accuracy   : {metrics.get('accuracy', 'n/a'):.2%}",
            f"ROC‑AUC    : {metrics.get('roc_auc', 'n/a'):.3f}",
            f"Precision  : {metrics.get('precision', 'n/a'):.2%}",
            f"Recall     : {metrics.get('recall', 'n/a'):.2%}",
            f"F1‑Score   : {metrics.get('f1', 'n/a'):.2%}",
        ]
        add_bullet_slide(prs, "Classification Engine – Evaluation Metrics", metric_bullets)
    else:
        add_bullet_slide(prs, "Classification Engine – Evaluation Metrics", ["Metrics file not found. Run classification_engine.py first."])

    # 4️⃣ Recommendation Engine – Sample data table
    if RECOMMENDATIONS_PATH.exists():
        df_rec = pd.read_csv(RECOMMENDATIONS_PATH)
        add_table_slide(prs, "Recommendation Snapshot (Top 8 rows)", df_rec)
    else:
        add_bullet_slide(prs, "Recommendation Engine", ["Recommendation CSV not found. Ensure dummy dataset exists."])

    # 5️⃣ Next Steps
    next_steps = [
        "Fine‑tune classifier hyper‑parameters and evaluate on hold‑out set.",
        "Enrich recommendation data with real product attributes and propensity scores.",
        "Integrate both engines into a real‑time decision service.",
        "A/B test recommendations to measure uplift in conversion rate.",
    ]
    add_bullet_slide(prs, "Next Steps & Recommendations", next_steps)

    # Save the deck
    prs.save(OUTPUT_PPTX)
    print(f"Presentation generated at {OUTPUT_PPTX}")

if __name__ == "__main__":
    main()
